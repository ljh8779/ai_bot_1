from io import BytesIO

import pypdfium2
import pytesseract
from PIL import Image, ImageFilter, ImageOps
from pypdf import PdfReader
from pptx import Presentation

from app.config import get_settings

settings = get_settings()
IMAGE_UPLOAD_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}
SUPPORTED_UPLOAD_SUFFIXES = {".txt", ".md", ".pdf", ".pptx", *IMAGE_UPLOAD_SUFFIXES}


def _decode_plain_text(raw: bytes) -> str:
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("cp949", errors="ignore")


def _content_score(text: str) -> int:
    # Normalize whitespace for rough quality scoring.
    return len(" ".join((text or "").split()))


def _preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    processed = image.convert("L")
    processed = ImageOps.autocontrast(processed)
    processed = processed.filter(ImageFilter.MedianFilter(size=3))

    if min(processed.size) < 1400:
        processed = processed.resize((processed.width * 2, processed.height * 2), Image.Resampling.LANCZOS)

    processed = processed.point(lambda p: 255 if p > 150 else 0)
    return processed


def _ocr_with_best_config(image: Image.Image) -> str:
    configs = [
        "--oem 1 --psm 6",
        "--oem 1 --psm 11",
    ]
    best_text = ""
    best_score = -1

    for config in configs:
        text = pytesseract.image_to_string(image, lang=settings.ocr_tesseract_lang, config=config)
        score = _content_score(text)
        if score > best_score:
            best_score = score
            best_text = text

    return best_text.strip()


def _extract_pdf_text(raw: bytes) -> str:
    reader = PdfReader(BytesIO(raw))
    lines: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            lines.append(text.strip())
    return "\n\n".join(lines)


def _extract_pdf_text_with_ocr(raw: bytes) -> str:
    scale = max(1.0, settings.pdf_ocr_dpi / 72.0)
    lines: list[str] = []

    pdf = pypdfium2.PdfDocument(BytesIO(raw))
    try:
        page_count = len(pdf)
        page_limit = min(page_count, settings.pdf_ocr_max_pages)
        for index in range(page_limit):
            page = pdf[index]
            bitmap = page.render(scale=scale)
            try:
                image = bitmap.to_pil().convert("RGB")
                text = _ocr_with_best_config(_preprocess_image_for_ocr(image))
                if text and text.strip():
                    lines.append(text.strip())
            finally:
                bitmap.close()
                page.close()
    finally:
        pdf.close()
    return "\n\n".join(lines)


def _extract_image_text_with_ocr(raw: bytes) -> str:
    with Image.open(BytesIO(raw)) as image:
        text = _ocr_with_best_config(_preprocess_image_for_ocr(image.convert("RGB")))
    return text.strip()


def _extract_pptx_text(raw: bytes) -> str:
    presentation = Presentation(BytesIO(raw))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
    return "\n".join(lines)


def _sanitize_text(content: str) -> str:
    normalized = content.strip()
    return normalized.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore").replace("\x00", " ").strip()


def extract_content_from_upload(suffix: str, raw: bytes) -> str:
    if suffix not in SUPPORTED_UPLOAD_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_UPLOAD_SUFFIXES))
        raise ValueError(f"Unsupported file type: {suffix}. Supported: {supported}")

    ocr_error: str | None = None
    content = ""

    if suffix in {".txt", ".md"}:
        content = _decode_plain_text(raw)
    elif suffix == ".pdf":
        content = _extract_pdf_text(raw)
        if settings.pdf_ocr_enabled and _content_score(content) < settings.pdf_ocr_fallback_min_chars:
            try:
                ocr_text = _extract_pdf_text_with_ocr(raw)
                if _content_score(ocr_text) > _content_score(content):
                    content = ocr_text
            except Exception as exc:
                ocr_error = str(exc)
    elif suffix in IMAGE_UPLOAD_SUFFIXES:
        try:
            content = _extract_image_text_with_ocr(raw)
        except Exception as exc:
            ocr_error = str(exc)
    else:
        content = _extract_pptx_text(raw)

    safe_text = _sanitize_text(content)
    if safe_text:
        return safe_text

    message = "No extractable text found in the uploaded file."
    if suffix == ".pdf":
        message += " The PDF may be image-only. OCR fallback was attempted."
    if suffix in IMAGE_UPLOAD_SUFFIXES:
        message += " OCR could not extract readable text from the image."
    if ocr_error:
        message += f" OCR detail: {ocr_error}"
    raise ValueError(message)
